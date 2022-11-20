import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from generate_plot import generate_plot
import json

data = json.load(open("data.json"))
pres = Presentation('slide_plan.pptx')


def first():
    slide = pres.slides[0]

    # title
    slide.shapes.title.text = data["first_slide"]["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # subtitle
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left * 2, top * 1.5, width * 10, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = data["first_slide"]["subtitle"]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)

    # image
    top_img_path = data["first_slide"]["image_path"]

    top_pic = slide.shapes.add_picture(
        top_img_path,
        left=Inches(3.5),
        top=Inches(3.3),
        width=Inches(3)
    )

    ref_element = slide.shapes[0]._element
    ref_element.addnext(top_pic._element)


def second():
    slide = pres.slides[1]

    slide.shapes[0].text_frame.text = data["second_slide"]["text_1_1"]
    slide.shapes[1].text_frame.text = data["second_slide"]["text_1_2"]

    slide.shapes[2].text_frame.text = data["second_slide"]["text_2_1"]
    slide.shapes[1].text_frame.paragraphs[0].font.size = Pt(16)

    slide.shapes[3].text_frame.text = data["second_slide"]["text_2_2"]
    slide.shapes[3].text_frame.paragraphs[0].font.size = Pt(16)


def third():
    slide = pres.slides[2]
    # title
    slide.shapes[1].text = data["third_slide"]["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    # text
    slide.shapes[4].text = data["third_slide"]["paragraph_1"] + data["third_slide"]["text_1"] +\
                           data["third_slide"]["paragraph_2"] + data["third_slide"]["text_2"] +\
                           data["third_slide"]["paragraph_3"] + data["third_slide"]["text_3"]
    for paragraph in slide.shapes[4].text_frame.paragraphs:
        paragraph.font.size = Pt(16)

def fourth():

    def add_paragraph(text):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

    slide = pres.slides[3]
    next_year, current_year = sorted(data["fourth_slide"]["years"], key=int, reverse=True)[:2]

    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left*0.2, top * 3, width*2, height*4)
    tf = txBox.text_frame
    # выручка текущего года
    add_paragraph(f"{data['fourth_slide']['paragraph_1']}({current_year}) - {data['fourth_slide']['years'][str(current_year)][data['fourth_slide']['paragraph_1']]}")
    # прибыль текущего года
    add_paragraph(f"{data['fourth_slide']['paragraph_2']}({current_year}) - {data['fourth_slide']['years'][str(current_year)][data['fourth_slide']['paragraph_2']]}")
    # выручка следующего года
    add_paragraph(f"{data['fourth_slide']['paragraph_3']}({next_year}) - {data['fourth_slide']['years'][str(next_year)][data['fourth_slide']['paragraph_1']]}")
    # прибыль следующего года
    add_paragraph(f"{data['fourth_slide']['paragraph_4']}({next_year}) - {data['fourth_slide']['years'][str(next_year)][data['fourth_slide']['paragraph_2']]}")

    # plot
    generate_plot()
    top_pic = slide.shapes.add_picture(
        data["fourth_slide"]["path2plot"],
        left=Inches(4),
        top=Inches(2),
        width=Inches(6)
    )

    ref_element = slide.shapes[0]._element
    ref_element.addnext(top_pic._element)

def fifth():
    slide = pres.slides[4]

    # first text block
    slide.shapes[3].text = data["fifth_slide"]["text_1"]
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # second text block
    slide.shapes[4].text = data["fifth_slide"]["text_2"]
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

def sixth():
    slide = pres.slides[5]
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left * 2, top * 1.7, width * 10, height)
    tf = txBox.text_frame

    for name, mem_data in data["sixth_slide"]["members"].items():
        p = tf.add_paragraph()
        p.text = f"• {name}\n\tposition: {mem_data['position']}\n\texperience: {mem_data['experience']}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0,0,0)

def seventh():
    slide = pres.slides[6]
    slide.shapes[4].text = f"ФИО: {data['seventh_slide']['name']}\n" +\
                            f"Номер телефона: {data['seventh_slide']['number']}\n" + \
                           f"Электронная почта: {data['seventh_slide']['email']}\n" + \
                           f"Телеграмм: {data['seventh_slide']['telegram']}\n"
    for paragraph in slide.shapes[4].text_frame.paragraphs:
        paragraph.font.size = Pt(16)

if __name__ == "__main__":
    first()
    second()
    third()
    fourth()
    fifth()
    sixth()
    seventh()
    pres.save('slide_test.pptx')
