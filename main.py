import collections.abc
from pptx import Presentation
from pptx.util import Inches

from PIL import Image
import numpy as np
import cv2

# convert pixels to inches
def px_to_inches(path):
    im = Image.open(path)
    width = im.width
    height = im.height

    return (width, height)

def prepare_background(
        img_path,
        background_size = (325, 720, 3),
        split_coeff = 1.3
):
    image = cv2.imread(img_path)
    background = np.ones(background_size)*255

    while image.shape[0] > background_size[0] or image.shape[1] > background_size[1]:
        image = cv2.resize(image, [x//split_coeff for x in image.shape[:2][::-1]])

    print(image.shape, background.shape)
    margin_top = (background_size[0] - image.shape[0]) // 2
    margin_left = (background_size[1] - image.shape[1]) // 2

    y_min = margin_top
    x_min = margin_left
    y_max = margin_top + image.shape[0]
    x_max = margin_left + image.shape[1]

    background[y_min:y_max, x_min:x_max] = image

    cv2.imwrite("first_slide.jpg", background)

pres = Presentation()


def first_slide():

    title_slide_layout = pres.slide_layouts[1]
    slide = pres.slides.add_slide(title_slide_layout)
    left = top = Inches(0)
    img_path = 'background.jpg'
    pic = slide.shapes.add_picture(img_path, left, top, width=pres.slide_width, height=pres.slide_height)

    slide.shapes.title.text = " Умные часы"
    slide.placeholders[1].text = " С использованием машинного обучения"

    # This moves it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


    # change slide sizes to Widescreen
    prepare_background('mibands.jpg')
    slide.shapes.add_picture('background_banner.jpg', left=Inches(0), top=Inches(3))

def second_slide():
    pass

def third_slide():
    pass

def fourth_slide():
    pass

def fifth_slide():
    pass

def sixth_slide():
    pass

first_slide()

pres.save('slide.pptx')